"""
loader.py
─────────
Loads content from every file type into a unified list of pages.
Each page = { "text": str, "source": str, "page": int, "type": str }

Supported: PDF, DOCX, PPTX, TXT, images (JPG/PNG via LLM OCR)
"""

import os
import base64
import json
from config import groq_client as client, LLM_MODEL as MODEL_NAME


# ─────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────

def _encode_image(path):
    with open(path, "rb") as f:
        return base64.b64encode(f.read()).decode()


# ─────────────────────────────────────────────────────────────────
# LOADERS
# ─────────────────────────────────────────────────────────────────

def load_pdf(path):
    """Extract text page by page from a PDF."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            if text:
                pages.append({
                    "text":   text,
                    "source": os.path.basename(path),
                    "page":   i + 1,
                    "type":   "pdf"
                })
        doc.close()
        return pages
    except Exception as e:
        print(f"  ⚠️ PDF load failed ({path}): {e}")
        return []


def load_docx(path):
    """Extract paragraphs from a Word document."""
    try:
        from docx import Document
        doc = Document(path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if text:
            return [{"text": text, "source": os.path.basename(path), "page": 1, "type": "docx"}]
        return []
    except Exception as e:
        print(f"  ⚠️ DOCX load failed ({path}): {e}")
        return []


def load_pptx(path):
    """Extract text slide by slide from a PowerPoint file."""
    try:
        from pptx import Presentation
        prs   = Presentation(path)
        pages = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                pages.append({
                    "text":   "\n".join(texts),
                    "source": os.path.basename(path),
                    "page":   i + 1,
                    "type":   "pptx"
                })
        return pages
    except Exception as e:
        print(f"  ⚠️ PPTX load failed ({path}): {e}")
        return []


def load_txt(path):
    """Load plain text file."""
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read().strip()
        if text:
            return [{"text": text, "source": os.path.basename(path), "page": 1, "type": "txt"}]
        return []
    except Exception as e:
        print(f"  ⚠️ TXT load failed ({path}): {e}")
        return []


def load_image(path):
    """Use LLM OCR to extract text from a handwritten or printed image."""
    try:
        encoded = _encode_image(path)
        prompt  = """Extract ALL text from this image exactly as written.
Return plain text only — no JSON, no markdown, no extra commentary.
Preserve paragraph structure with newlines."""

        res = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": [
                {"type": "text",      "text": prompt},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{encoded}"}}
            ]}],
            temperature=0.0
        )
        text = res.choices[0].message.content.strip()
        if text:
            return [{"text": text, "source": os.path.basename(path), "page": 1, "type": "image"}]
        return []
    except Exception as e:
        print(f"  ⚠️ Image OCR failed ({path}): {e}")
        return []


# ─────────────────────────────────────────────────────────────────
# MAIN ENTRY
# ─────────────────────────────────────────────────────────────────

SUPPORTED = {
    ".pdf":  load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".txt":  load_txt,
    ".jpg":  load_image,
    ".jpeg": load_image,
    ".png":  load_image,
}


def load_file(path):
    """Load a single file. Returns list of page dicts."""
    ext = os.path.splitext(path)[1].lower()
    loader = SUPPORTED.get(ext)
    if not loader:
        print(f"  ⚠️ Unsupported file type: {ext}")
        return []
    print(f"  📄 Loading: {os.path.basename(path)} ({ext})")
    return loader(path)


def load_folder(folder_path):
    """
    Load all supported files from a folder.
    Returns flat list of all pages across all files.
    """
    all_pages = []
    for fname in sorted(os.listdir(folder_path)):
        ext = os.path.splitext(fname)[1].lower()
        if ext in SUPPORTED:
            fpath = os.path.join(folder_path, fname)
            pages = load_file(fpath)
            all_pages.extend(pages)
            print(f"    → {len(pages)} page(s) loaded")
    print(f"\n  ✅ Total pages loaded: {len(all_pages)}")
    return all_pages